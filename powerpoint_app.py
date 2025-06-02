import streamlit as st
import openai
import os
import pptx
from pptx import Presentation
from pptx.util import Inches
import json
from typing import List, Dict
import hashlib
from datetime import datetime
from openai import AsyncOpenAI
import asyncio
import tempfile
import io

def powerpoint_app():
    # Titel der App
    st.title("BonsAI PowerPoint Sprachprüfung und Korrektur")

    # --- Constants and Configurations ---

    SUPPORTED_EXTENSIONS = ('.pptx',)

    # Model options for the dropdown
    MODEL_OPTIONS = {
        "GPT-4.1-mini": "gpt-4.1-mini",
        "GPT-4o": "gpt-4o"
    }

    # Language options for the dropdown
    LANGUAGE_OPTIONS = {
        "Deutsch": "de",
        "Spanisch": "es",
        "Französisch": "fr", 
        "Englisch": "en",
        "Italienisch": "it",
        "Portugiesisch": "pt",
        "Chinesisch (Vereinfacht)": "zh-CN",
        "Chinesisch (Traditionell)": "zh-TW",
        "Japanisch": "ja",
        "Koreanisch": "ko",
        "Russisch": "ru",
        "Arabisch": "ar",
        "Hindi": "hi",
        "Niederländisch": "nl",
        "Schwedisch": "sv",
        "Norwegisch": "no",
        "Dänisch": "da",
        "Finnisch": "fi",
        "Polnisch": "pl",
        "Tschechisch": "cs",
        "Ungarisch": "hu",
        "Rumänisch": "ro",
        "Bulgarisch": "bg",
        "Kroatisch": "hr",
        "Serbisch": "sr",
        "Slowakisch": "sk",
        "Slowenisch": "sl",
        "Estnisch": "et",
        "Lettisch": "lv",
        "Litauisch": "lt",
        "Griechisch": "el",
        "Türkisch": "tr",
        "Hebräisch": "he",
        "Thai": "th",
        "Vietnamesisch": "vi",
        "Indonesisch": "id",
        "Malaiisch": "ms",
        "Filipino": "fil",
        "Ukrainisch": "uk"
    }

    # Default system prompt
    DEFAULT_SYSTEM_PROMPT = """Du bist ein hilfreicher Assistent, der Texte in {target_language} übersetzt.
Behalte die ursprüngliche Bedeutung so genau wie möglich bei.
Passe den Ton der Übersetzung so an, dass er für professionelle Präsentationen in der Zielsprache ({target_language}) angemessen ist.
Der übersetzte Text sollte ungefähr die gleiche Zeichenlänge wie der ursprüngliche Text haben (innerhalb einer 5%-Marge).
Übersetze keine E-Mails, Telefonnummern oder andere nicht-textuelle Inhalte.
Verwende korrekte Umlaute und Sonderzeichen für die Zielsprache.
Gib die Übersetzung als JSON-Objekt genau wie folgt zurück: {{"translated": "<übersetzter Text>"}}"""

    # --- Helper Functions ---

    def generate_prompt_hash(prompt: str) -> str:
        """Generates a SHA-256 hash of the prompt for use as a cache key."""
        return hashlib.sha256(prompt.encode('utf-8')).hexdigest()

    def safe_text_extraction(text: str) -> str:
        """Safely extracts and normalizes text to handle encoding issues."""
        if not text:
            return ""
        
        # Ensure proper UTF-8 encoding
        try:
            # If text is bytes, decode it
            if isinstance(text, bytes):
                text = text.decode('utf-8', errors='replace')
            
            # Normalize the text to handle any encoding issues
            text = text.encode('utf-8', errors='replace').decode('utf-8')
            
            # Clean up any problematic characters while preserving umlauts
            text = text.strip()
            
            return text
        except Exception as e:
            st.warning(f"Textverarbeitungsfehler: {e}")
            return str(text) if text else ""

    async def translate_text_with_openai(prompt: str, target_language: str, cache: Dict, model: str = "gpt-4.1-mini", system_prompt: str = None, max_retries: int = 3) -> str:
        """Translates text using the OpenAI API, with caching and retries."""
        # Ensure proper text encoding
        prompt = safe_text_extraction(prompt)
        
        # Include system prompt in hash for separate caching
        cache_key = prompt + model + (system_prompt or DEFAULT_SYSTEM_PROMPT)
        prompt_hash = generate_prompt_hash(cache_key)
        if prompt_hash in cache:
            return cache[prompt_hash]

        # Use custom system prompt or default
        if system_prompt is None:
            system_instruction = DEFAULT_SYSTEM_PROMPT.format(target_language=target_language)
        else:
            system_instruction = system_prompt.format(target_language=target_language)

        api_key = st.session_state.get("api_key")
        if not api_key:
            raise ValueError("OpenAI API-Schlüssel nicht gefunden. Bitte gib deinen API-Schlüssel ein.")

        client = AsyncOpenAI(api_key=api_key, timeout=30.0)

        for attempt in range(max_retries):
            try:
                response = await client.chat.completions.create(
                    model=model,
                    messages=[
                        {"role": "system", "content": system_instruction},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.2,
                    max_tokens=4096,
                    timeout=30,
                    response_format={"type": "json_object"}
                )
                result = response.choices[0].message.content.strip()

                for parse_attempt in range(max_retries):
                    try:
                        parsed = json.loads(result)
                        translated_text = parsed["translated"]
                        # Ensure proper encoding of the translated text
                        translated_text = safe_text_extraction(translated_text)
                        cache[prompt_hash] = translated_text
                        return translated_text
                    except (json.JSONDecodeError, KeyError) as e:
                        if parse_attempt == max_retries - 1:
                            return prompt
                    except Exception as e:
                        return prompt

            except Exception as e:
                if attempt == max_retries - 1:
                    return prompt
        return prompt

    def extract_text_from_presentation(presentation_path: str) -> List[Dict]:
        """Extracts text and context from a PowerPoint presentation."""
        try:
            prs = Presentation(presentation_path)
            text_data = []

            for slide_number, slide in enumerate(prs.slides, start=1):
                for shape_index, shape in enumerate(slide.shapes):
                    shape_id = f"slide{slide_number}_shape{shape_index}"
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    # Safely extract and normalize text
                                    clean_text = safe_text_extraction(run.text)
                                    if clean_text:
                                        shape_type = "UNKNOWN"
                                        if shape == slide.shapes.title:
                                            shape_type = "TITLE"
                                        elif shape.has_table:
                                            shape_type = "TABLE"
                                        else:
                                            shape_type = "BODY"

                                        text_data.append({
                                            "slide_number": slide_number,
                                            "shape_type": shape_type,
                                            "text": clean_text,
                                            "shape_id": shape_id,
                                        })

                    elif shape.has_table:
                        for row_idx, row in enumerate(shape.table.rows):
                            for col_idx, cell in enumerate(row.cells):
                                if cell.text.strip():
                                    # Safely extract and normalize text
                                    clean_text = safe_text_extraction(cell.text)
                                    if clean_text:
                                        text_data.append({
                                            "slide_number": slide_number,
                                            "shape_type": "TABLE",
                                            "text": clean_text,
                                            "shape_id": f"{shape_id}_row{row_idx}_col{col_idx}"
                                        })

            return text_data

        except Exception as e:
            st.error(f"Fehler beim Extrahieren von Text aus der Präsentation: {e}")
            return []

    async def batch_translate_texts_with_openai(text_entries: List[Dict], target_language: str, cache: Dict, model: str = "gpt-4.1-mini", system_prompt: str = None, max_retries: int = 3, batch_size: int = 10) -> None:
        """Batch translates multiple texts using the OpenAI API with structured JSON output."""
        texts_to_translate = []
        cache_key_base = model + (system_prompt or DEFAULT_SYSTEM_PROMPT)
        
        for entry in text_entries:
            # Ensure proper text encoding
            clean_text = safe_text_extraction(entry["text"])
            cache_key = clean_text + cache_key_base
            prompt_hash = generate_prompt_hash(cache_key)
            if prompt_hash not in cache:
                texts_to_translate.append((prompt_hash, clean_text))

        if not texts_to_translate:
            return

        total_batches = (len(texts_to_translate) + batch_size - 1) // batch_size
        
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        tasks = []
        for batch_num, i in enumerate(range(0, len(texts_to_translate), batch_size), 1):
            batch = texts_to_translate[i:i + batch_size]
            payload = {hash_: text for hash_, text in batch}

            status_text.text(f"Verarbeite Batch {batch_num}/{total_batches} ({len(batch)} Texte)")
            progress_bar.progress(batch_num / total_batches)

            # Use custom system prompt or default for batch translation
            if system_prompt is None:
                system_instruction = f"""Du bist ein hilfreicher Assistent, der mehrere Texte in {target_language} übersetzt.
Behalte die ursprüngliche Bedeutung so genau wie möglich bei.
Passe den Ton jeder Übersetzung so an, dass er für professionelle Präsentationen in der Zielsprache ({target_language}) angemessen ist.
Der übersetzte Text für jede Eingabe sollte ungefähr die gleiche Länge wie der ursprüngliche Text haben (innerhalb einer 10%-Marge).
Verwende korrekte Umlaute und Sonderzeichen für die Zielsprache.
Gib die Übersetzungen als JSON-Objekt genau wie folgt zurück:
{{"translations": {{"<sha256 hash>": "<übersetzter Text>"}} }}"""
            else:
                system_instruction = system_prompt.format(target_language=target_language) + f"""
Verwende korrekte Umlaute und Sonderzeichen für die Zielsprache.
Gib die Übersetzungen als JSON-Objekt genau wie folgt zurück:
{{"translations": {{"<sha256 hash>": "<übersetzter Text>"}} }}"""

            prompt_data = {
                "texts": payload,
                "target_language": target_language,
                "instructions": "Translate each text, maintaining original meaning and formatting. Use correct umlauts and special characters."
            }

            api_key = st.session_state.get("api_key")
            if not api_key:
                raise ValueError("OpenAI API-Schlüssel nicht gefunden. Bitte gib deinen API-Schlüssel ein.")

            client = AsyncOpenAI(api_key=api_key, timeout=60.0)
            tasks.append(translate_batch(client, system_instruction, prompt_data, cache, max_retries, batch_num, total_batches, model))

        await asyncio.gather(*tasks)
        progress_bar.progress(1.0)
        status_text.text("Übersetzung abgeschlossen!")

    async def translate_batch(client: AsyncOpenAI, system_instruction: str, prompt_data: Dict, cache: Dict, max_retries: int, batch_num: int, total_batches: int, model: str) -> None:
        """Translates a single batch (async)."""
        for attempt in range(max_retries):
            try:
                response = await client.chat.completions.create(
                    model=model,
                    messages=[
                        {"role": "system", "content": system_instruction},
                        {"role": "user", "content": json.dumps(prompt_data, ensure_ascii=False)}
                    ],
                    temperature=0.2,
                    max_tokens=4096,
                    timeout=60,
                    response_format={"type": "json_object"}
                )
                output = response.choices[0].message.content.strip()

                for parse_attempt in range(max_retries):
                    try:
                        result = json.loads(output)
                        translations = result.get("translations", {})
                        for hash_, translated_text in translations.items():
                            # Ensure proper encoding of translated text
                            clean_translated_text = safe_text_extraction(translated_text)
                            cache[hash_] = clean_translated_text
                        break
                    except (json.JSONDecodeError, KeyError) as e:
                        if parse_attempt == max_retries - 1:
                            pass
                    except Exception as e:
                        pass
                else:
                    continue
                break

            except Exception as e:
                if attempt == max_retries - 1:
                    pass

    async def translate_presentation(presentation_file, target_language: str, model: str = "gpt-4.1-mini", system_prompt: str = None) -> bytes:
        """Translates a PowerPoint presentation and returns the translated version as bytes."""
        
        # Create temporary files with proper encoding
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_input:
            temp_input.write(presentation_file.read())
            temp_input_path = temp_input.name
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_output:
            temp_output_path = temp_output.name

        try:
            text_data = extract_text_from_presentation(temp_input_path)
            if not text_data:
                st.warning("Kein Text zum Übersetzen in der Präsentation gefunden.")
                return None

            # Initialize cache for this session
            cache = {}

            await batch_translate_texts_with_openai(text_data, target_language, cache, model, system_prompt)

            # Load and save the presentation
            prs = Presentation(temp_input_path)
            prs.save(temp_output_path)

            translated_prs = Presentation(temp_output_path)

            # Apply translations
            translated_text_data = []
            cache_key_base = model + (system_prompt or DEFAULT_SYSTEM_PROMPT)
            
            for text_entry in text_data:
                clean_text = safe_text_extraction(text_entry["text"])
                cache_key = clean_text + cache_key_base
                prompt_hash = generate_prompt_hash(cache_key)
                translated_text = cache.get(prompt_hash, clean_text)
                translated_text_entry = text_entry.copy()
                translated_text_entry["translated_text"] = translated_text
                translated_text_data.append(translated_text_entry)

            for slide_number, slide in enumerate(translated_prs.slides, start=1):
                for shape_index, shape in enumerate(slide.shapes):
                    shape_id = f"slide{slide_number}_shape{shape_index}"

                    if shape.has_text_frame:
                        try:
                            text_frame = shape.text_frame
                            for paragraph in text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text.strip():
                                        clean_original_text = safe_text_extraction(run.text)
                                        for entry in translated_text_data:
                                            if entry["shape_id"] == shape_id and safe_text_extraction(entry["text"]) == clean_original_text:
                                                # Ensure proper encoding when setting the text
                                                run.text = entry["translated_text"]
                                                break
                        except Exception as e:
                            continue

                    elif shape.has_table:
                        try:
                            for row_idx, row in enumerate(shape.table.rows):
                                for col_idx, cell in enumerate(row.cells):
                                    cell_shape_id = f"{shape_id}_row{row_idx}_col{col_idx}"
                                    cell_translated_text_entry = next((entry for entry in translated_text_data if entry["shape_id"] == cell_shape_id), None)
                                    if cell_translated_text_entry:
                                        # Ensure proper encoding when setting the text
                                        cell.text = cell_translated_text_entry["translated_text"]
                        except Exception as e:
                            continue

            translated_prs.save(temp_output_path)
            
            # Read the translated file as bytes
            with open(temp_output_path, 'rb') as f:
                translated_bytes = f.read()
            
            return translated_bytes

        except Exception as e:
            st.error(f"Fehler während des Übersetzungsprozesses: {e}")
            return None
        finally:
            # Clean up temporary files
            try:
                os.unlink(temp_input_path)
                os.unlink(temp_output_path)
            except:
                pass

    # Main Streamlit app content
    st.markdown("Übersetze deine PowerPoint-Präsentationen mit OpenAI's GPT-Modellen")
    
    # Sidebar for configuration
    with st.sidebar:
        st.header("Konfiguration")
        
        # API Key input
        api_key = st.text_input(
            "OpenAI API-Schlüssel",
            type="password",
            help="Gib deinen OpenAI API-Schlüssel ein, um den Übersetzungsservice zu nutzen"
        )
        
        if api_key:
            st.session_state["api_key"] = api_key
            st.success("✅ API-Schlüssel gesetzt!")
        
        # Model selection
        selected_model_name = st.selectbox(
            "KI-Modell",
            options=list(MODEL_OPTIONS.keys()),
            help="Wähle das KI-Modell für die Übersetzung. GPT-4.1-mini ist schneller und günstiger, GPT-4o bietet höchste Qualität."
        )
        
        selected_model = MODEL_OPTIONS[selected_model_name]
        
        # Show model info
        if "mini" in selected_model:
            st.info("💡 GPT-4.1-mini: Schneller & 83% günstiger als GPT-4o")
        else:
            st.info("🎯 GPT-4o: Höchste Qualität & Genauigkeit")
        
        # Language selection
        selected_language_name = st.selectbox(
            "Zielsprache",
            options=list(LANGUAGE_OPTIONS.keys()),
            help="Wähle die Sprache aus, in die du deine Präsentation übersetzen möchtest"
        )
        
        target_language = LANGUAGE_OPTIONS[selected_language_name]
        
        st.info(f"Ausgewählt: {selected_language_name} ({target_language})")
    
    # System prompt customization (collapsed by default)
    with st.expander("⚙️ Systemprompt anpassen (Erweitert)", expanded=False):
        st.markdown("**Hier kannst du das Systemprompt für die Übersetzung anpassen:**")
        
        custom_system_prompt = st.text_area(
            "Systemprompt",
            value=DEFAULT_SYSTEM_PROMPT,
            height=150,
            help="Verwende {target_language} als Platzhalter für die Zielsprache. Das Prompt sollte Anweisungen für JSON-Ausgabe enthalten."
        )
        
        if st.button("🔄 Standard wiederherstellen"):
            st.rerun()
        
        # Show preview of formatted prompt
        if target_language:
            st.markdown("**Vorschau (formatiert):**")
            preview = custom_system_prompt.format(target_language=selected_language_name)
            st.code(preview, language="text")
    
    # Main content area
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.header("Präsentation hochladen")
        
        uploaded_file = st.file_uploader(
            "PowerPoint-Datei auswählen",
            type=['pptx'],
            help="Lade eine .pptx-Datei zum Übersetzen hoch"
        )
        
        if uploaded_file is not None:
            st.success(f"✅ Datei hochgeladen: {uploaded_file.name}")
            
            # Display file info
            file_size = len(uploaded_file.getvalue()) / 1024 / 1024  # MB
            st.info(f"Dateigröße: {file_size:.2f} MB")
    
    with col2:
        st.header("Übersetzung")
        
        if uploaded_file is not None and api_key:
            if st.button("🚀 Präsentation übersetzen", type="primary"):
                with st.spinner("Präsentation wird übersetzt..."):
                    try:
                        # Reset file pointer
                        uploaded_file.seek(0)
                        
                        # Use custom system prompt if different from default
                        system_prompt_to_use = custom_system_prompt if custom_system_prompt != DEFAULT_SYSTEM_PROMPT else None
                        
                        # Translate the presentation
                        translated_bytes = asyncio.run(
                            translate_presentation(uploaded_file, target_language, selected_model, system_prompt_to_use)
                        )
                        
                        if translated_bytes:
                            # Generate download filename
                            original_name = uploaded_file.name.replace('.pptx', '')
                            model_suffix = "mini" if "mini" in selected_model else "4o"
                            download_filename = f"{original_name}_übersetzt_{target_language}_{model_suffix}.pptx"
                            
                            st.success("🎉 Übersetzung abgeschlossen!")
                            
                            # Download button
                            st.download_button(
                                label="📥 Übersetzte Präsentation herunterladen",
                                data=translated_bytes,
                                file_name=download_filename,
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        else:
                            st.error("Übersetzung fehlgeschlagen. Bitte versuche es erneut.")
                            
                    except Exception as e:
                        st.error(f"Ein Fehler ist aufgetreten: {str(e)}")
        
        elif not api_key:
            st.warning("⚠️ Bitte gib deinen OpenAI API-Schlüssel in der Seitenleiste ein")
        elif uploaded_file is None:
            st.info("📤 Bitte lade eine PowerPoint-Datei hoch, um zu beginnen")
    
    # Model comparison info
    with st.expander("🔍 Modell-Vergleich"):
        st.markdown("""
        | Modell | Geschwindigkeit | Kosten | Qualität | Beste Verwendung |
        |--------|----------------|--------|----------|------------------|
        | **GPT-4.1-mini** | Schneller | 83% günstiger | Sehr gut | Alltägliche Übersetzungen, große Mengen |
        | **GPT-4o** | Standard | Standard | Höchste | Wichtige Dokumente, maximale Genauigkeit |
        
        **GPT-4.1-mini Vorteile:**
        - ⚡ Deutlich schnellere Verarbeitung
        - 💰 Erheblich niedrigere Kosten
        - 🎯 Sehr gute Qualität für die meisten Anwendungsfälle
        - 📄 1 Million Token Kontext (wie GPT-4o)
        """)
    
    # Instructions
    with st.expander("📖 Wie man diese App verwendet"):
        st.markdown("""
        1. **OpenAI API-Schlüssel besorgen**: Frag Tobias oder Jonathan um den API-Schlüssel zu erhalten
        2. **API-Schlüssel eingeben**: Füge deinen API-Schlüssel in der Seitenleiste ein (er wird sicher in deiner Sitzung gespeichert)
        3. **Modell auswählen**: Wähle zwischen GPT-4.1-mini (Standard, schneller & günstiger) oder GPT-4o (höchste Qualität)
        4. **Sprache auswählen**: Wähle deine Zielsprache aus dem Dropdown-Menü
        5. **Systemprompt anpassen** (optional): Passe das Übersetzungsverhalten im erweiterten Bereich an
        6. **Datei hochladen**: Lade deine PowerPoint (.pptx) Datei hoch
        7. **Übersetzen**: Klicke auf den Übersetzen-Button und warte, bis der Prozess abgeschlossen ist
        8. **Herunterladen**: Lade deine übersetzte Präsentation herunter
        
        **Hinweis**: Der Übersetzungsprozess kann je nach Größe deiner Präsentation einige Minuten dauern.
        """)
    
    # Footer
    st.markdown("---")

# Call the function when the script is run directly
if __name__ == "__main__":
    powerpoint_app() 